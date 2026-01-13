#!/usr/bin/env python3
"""
Bulk Document Ingestion Script (streaming + batching + guards + multi-format)
Processes documents from local directories, AWS S3 buckets, or Google Drive
and ingests them into a Pinecone namespace with streaming upserts.

Key protections:
- Skips empty/too-short chunks ("No valid texts provided" avoided)
- Embedding micro-batching (avoid per-request token limits)
- Optional page and chunk caps per file (tame huge PDFs/books)
- Broad format support with optional dependencies
- Legacy Office (.doc, .ppt) via LibreOffice conversion (fallbacks antiword/catppt)
"""

import os
import sys
import argparse
from pathlib import Path
from typing import List, Dict, Optional, Tuple
import time
import tempfile
import shutil
import subprocess
import re
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading

# Add parent directory to path to import services
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from services.pinecone_service import get_pinecone_service
from services.openai_service import get_openai_service
from services.bm25_service import get_bm25_service
from dotenv import load_dotenv

load_dotenv()

# ---------------- Optional deps (soft) ----------------
# PDF
try:
    import PyPDF2
    PDF_SUPPORT = True
except Exception:
    PDF_SUPPORT = False
    print("Warning: PyPDF2 not installed. PDF support disabled.")

# DOCX
try:
    import docx
    DOCX_SUPPORT = True
except Exception:
    DOCX_SUPPORT = False
    print("Warning: python-docx not installed. DOCX support disabled.")

# PPTX
try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except Exception:
    PPTX_SUPPORT = False
    print("Note: python-pptx not installed. PPTX support disabled.")

# XLSX
try:
    import openpyxl
    XLSX_SUPPORT = True
except Exception:
    XLSX_SUPPORT = False
    print("Note: openpyxl not installed. XLSX support disabled.")

# XLS legacy
try:
    import xlrd  # xlrd<2.0 supports xls
    XLS_SUPPORT = True
except Exception:
    XLS_SUPPORT = False
    print("Note: xlrd not installed (or no .xls support). XLS support disabled.")

# CSV/tabular helpers
try:
    import pandas as pd
    PANDAS_SUPPORT = True
except Exception:
    PANDAS_SUPPORT = False
    print("Note: pandas not installed. CSV/XLS/XLSX may use basic readers.")

# HTML
try:
    from bs4 import BeautifulSoup
    BS4_SUPPORT = True
except Exception:
    BS4_SUPPORT = False
    print("Note: beautifulsoup4 not installed. HTML parsing disabled.")

# RTF
try:
    from striprtf.striprtf import rtf_to_text
    RTF_SUPPORT = True
except Exception:
    RTF_SUPPORT = False
    print("Note: striprtf not installed. RTF support disabled.")

# ODT
try:
    from odf.opendocument import load as odf_load
    from odf.text import P
    ODT_SUPPORT = True
except Exception:
    ODT_SUPPORT = False
    print("Note: odfpy not installed. ODT support disabled.")

# EPUB
try:
    from ebooklib import epub
    EPUB_SUPPORT = True
except Exception:
    EPUB_SUPPORT = False
    print("Note: ebooklib not installed. EPUB support disabled.")

# Google Drive (lazy)
try:
    import importlib
    import io
    service_account = importlib.import_module('google.oauth2.service_account')
    gapi_discovery = importlib.import_module('googleapiclient.discovery')
    build = getattr(gapi_discovery, 'build')
    gapi_http = importlib.import_module('googleapiclient.http')
    MediaIoBaseDownload = getattr(gapi_http, 'MediaIoBaseDownload')
    GDRIVE_SUPPORT = True
except Exception:
    service_account = None
    build = None
    MediaIoBaseDownload = None
    import io as _io
    io = _io
    GDRIVE_SUPPORT = False

# S3
try:
    import boto3
    from botocore.exceptions import ClientError, NoCredentialsError
    S3_SUPPORT = True
except Exception:
    S3_SUPPORT = False

DEFAULT_CHUNK_SIZE = 1000
DEFAULT_CHUNK_OVERLAP = 200

ALLOWED_EXTENSIONS = {
    '.txt', '.md',
    '.pdf',
    '.docx', '.doc',    # <-- added .doc
    '.pages',           # <-- Apple Pages
    '.csv', '.xlsx', '.xls',
    '.pptx', '.ppt',    # <-- added .ppt
    '.html', '.htm',
    '.rtf',
    '.odt',
    '.epub',
}

# --------------- Chunking ---------------

def chunk_text(text: str, chunk_size: int, overlap: int) -> List[str]:
    if len(text) <= chunk_size:
        return [text.strip()]
    chunks, start = [], 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            for delim in ['. ', '! ', '? ', '\n\n']:
                last = text[start:end].rfind(delim)
                if last != -1:
                    end = start + last + len(delim)
                    break
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        start = max(start + 1, end - overlap)
    return chunks

# --------------- OCR helper ---------------

def _maybe_ocr_pdf(input_path: str, temp_dir: str) -> Optional[str]:
    try:
        subprocess.run(["ocrmypdf", "--version"], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
    except Exception:
        return None
    stem = Path(input_path).stem
    fd, out_path = tempfile.mkstemp(prefix=f"{stem}.ocr_", suffix=".pdf", dir=temp_dir)
    os.close(fd)
    try:
        if os.path.exists(out_path):
            os.unlink(out_path)
    except Exception:
        pass
    try:
        res = subprocess.run(
            ["ocrmypdf", "-l", "eng", "--optimize", "3", "--skip-text", input_path, out_path],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
        )
        if res.returncode == 0 and os.path.exists(out_path) and os.path.getsize(out_path) > 0:
            return out_path
        else:
            print("  OCR failed:", res.stderr.strip() or res.stdout.strip())
    except Exception as e:
        print(f"  OCR exception: {e}")
    try:
        if os.path.exists(out_path):
            os.unlink(out_path)
    except Exception:
        pass
    return None

# --------------- Legacy Office converters ---------------

def _which(cmd: str) -> Optional[str]:
    return shutil.which(cmd)

def _convert_with_soffice(input_path: str, out_ext: str, temp_dir: str) -> Optional[str]:
    """
    Convert using LibreOffice headless (soffice). Returns output path or None.
    out_ext: 'docx' or 'pptx' (we rely on our modern extractors afterward)
    """
    soffice = _which("soffice")
    if not soffice:
        return None
    try:
        # LibreOffice writes to --outdir with same basename but new ext
        res = subprocess.run(
            [soffice, "--headless", "--convert-to", out_ext, "--outdir", temp_dir, input_path],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
        )
        if res.returncode != 0:
            print("  soffice conversion failed:", res.stderr.strip() or res.stdout.strip())
            return None
        base = Path(input_path).stem + f".{out_ext}"
        candidate = os.path.join(temp_dir, base)
        if os.path.exists(candidate) and os.path.getsize(candidate) > 0:
            return candidate
    except Exception as e:
        print(f"  soffice conversion exception: {e}")
    return None

def _extract_with_antiword(input_path: str) -> Optional[str]:
    """Fallback for .doc → plain text if conversion unavailable."""
    antiword = _which("antiword")
    if not antiword:
        return None
    try:
        res = subprocess.run(
            [antiword, "-m", "UTF-8.txt", input_path],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
        )
        if res.returncode == 0 and res.stdout.strip():
            return res.stdout
        else:
            print("  antiword failed:", res.stderr.strip() or "no output")
    except Exception as e:
        print(f"  antiword exception: {e}")
    return None

def _extract_with_catppt(input_path: str) -> Optional[str]:
    """Fallback for .ppt → plain text if conversion unavailable."""
    catppt = _which("catppt")
    if not catppt:
        return None
    try:
        res = subprocess.run(
            [catppt, input_path],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
        )
        if res.returncode == 0 and res.stdout.strip():
            return res.stdout
        else:
            print("  catppt failed:", res.stderr.strip() or "no output")
    except Exception as e:
        print(f"  catppt exception: {e}")
    return None

# --------------- Extractors ---------------

def extract_text_from_pdf(file_path: str, max_pages: int = 0) -> Tuple[str, int]:
    if not PDF_SUPPORT:
        raise Exception("PDF support not installed")
    text, pages_read = "", 0
    with open(file_path, 'rb') as fh:
        reader = PyPDF2.PdfReader(fh)
        total = len(reader.pages)
        limit = total if max_pages <= 0 else min(total, max_pages)
        for i in range(limit):
            pages_read += 1
            page_text = reader.pages[i].extract_text() or ""
            text += page_text + "\n"
    return text, pages_read

def extract_text_from_docx(file_path: str) -> str:
    if not DOCX_SUPPORT:
        raise Exception("DOCX support not installed")
    d = docx.Document(file_path)
    return "\n".join(p.text for p in d.paragraphs)

def extract_text_from_pptx(file_path: str) -> str:
    if not PPTX_SUPPORT:
        raise Exception("PPTX support not installed")
    prs = Presentation(file_path)
    parts = []
    for idx, slide in enumerate(prs.slides, 1):
        parts.append(f"\n--- Slide {idx} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            parts.append("\n[Notes]")
            parts.append(slide.notes_slide.notes_text_frame.text or "")
    return "\n".join(parts)

def _join_cells(rows: List[List[Optional[str]]], max_rows: int) -> str:
    out = []
    for r_i, row in enumerate(rows):
        if r_i >= max_rows:
            out.append(f"... [truncated at {max_rows} rows]")
            break
        vals = []
        for v in row:
            s = "" if v is None else str(v)
            vals.append(s)
        out.append("\t".join(vals))
    return "\n".join(out)

def extract_text_from_xlsx(file_path: str, max_rows: int) -> str:
    if XLSX_SUPPORT:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"\n--- Sheet: {ws.title} ---")
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    break
                rows.append([cell for cell in row])
            parts.append(_join_cells(rows, max_rows))
        return "\n".join(parts)
    if PANDAS_SUPPORT:
        try:
            xl = pd.ExcelFile(file_path)
            parts = []
            for name in xl.sheet_names:
                df = xl.parse(name, nrows=max_rows)
                parts.append(f"\n--- Sheet: {name} ---")
                parts.append(df.to_csv(sep="\t", index=False))
            return "\n".join(parts)
        except Exception as e:
            raise Exception(f"Failed to read XLSX: {e}")
    raise Exception("XLSX support requires openpyxl or pandas")

def extract_text_from_xls(file_path: str, max_rows: int) -> str:
    if PANDAS_SUPPORT:
        try:
            xl = pd.ExcelFile(file_path)
            parts = []
            for name in xl.sheet_names:
                df = xl.parse(name, nrows=max_rows)
                parts.append(f"\n--- Sheet: {name} ---")
                parts.append(df.to_csv(sep="\t", index=False))
            return "\n".join(parts)
        except Exception as e:
            if XLS_SUPPORT:
                try:
                    book = xlrd.open_workbook(file_path)
                    parts = []
                    for si in range(book.nsheets):
                        sh = book.sheet_by_index(si)
                        parts.append(f"\n--- Sheet: {sh.name} ---")
                        rows = []
                        for r in range(min(max_rows, sh.nrows)):
                            row_vals = [sh.cell_value(r, c) for c in range(sh.ncols)]
                            rows.append(row_vals)
                        parts.append(_join_cells(rows, max_rows))
                    return "\n".join(parts)
                except Exception as e2:
                    raise Exception(f"Failed to read XLS via xlrd: {e2}")
            raise Exception(f"Failed to read XLS: {e}")
    if XLS_SUPPORT:
        try:
            book = xlrd.open_workbook(file_path)
            parts = []
            for si in range(book.nsheets):
                sh = book.sheet_by_index(si)
                parts.append(f"\n--- Sheet: {sh.name} ---")
                rows = []
                for r in range(min(max_rows, sh.nrows)):
                    row_vals = [sh.cell_value(r, c) for c in range(sh.ncols)]
                    rows.append(row_vals)
                parts.append(_join_cells(rows, max_rows))
            return "\n".join(parts)
        except Exception as e:
            raise Exception(f"Failed to read XLS: {e}")
    raise Exception("XLS support requires pandas (with suitable engine) or xlrd (legacy)")

def extract_text_from_csv(file_path: str, max_rows: int) -> str:
    if PANDAS_SUPPORT:
        try:
            df = pd.read_csv(file_path, nrows=max_rows)
            return df.to_csv(sep="\t", index=False)
        except Exception:
            pass
    out, count = [], 0
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        for line in f:
            out.append(line.rstrip("\n"))
            count += 1
            if count >= max_rows:
                out.append(f"... [truncated at {max_rows} rows]")
                break
    return "\n".join(out)

def extract_text_from_html(file_path: str, max_chars: int) -> str:
    if not BS4_SUPPORT:
        raise Exception("HTML support requires beautifulsoup4")
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    txt = soup.get_text(separator="\n", strip=True)
    if len(txt) > max_chars:
        return txt[:max_chars] + f"\n... [truncated at {max_chars} chars]"
    return txt

def extract_text_from_rtf(file_path: str) -> str:
    if not RTF_SUPPORT:
        raise Exception("RTF support requires striprtf")
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    return rtf_to_text(raw)

def extract_text_from_odt(file_path: str) -> str:
    if not ODT_SUPPORT:
        raise Exception("ODT support requires odfpy")
    doc = odf_load(file_path)
    texts = []
    for p in doc.getElementsByType(P):
        if p and p.firstChild:
            texts.append(str(p.firstChild))
    return "\n".join(texts)

def extract_text_from_epub(file_path: str) -> str:
    if not EPUB_SUPPORT or not BS4_SUPPORT:
        raise Exception("EPUB support requires ebooklib and beautifulsoup4")
    book = epub.read_epub(file_path)
    parts = []
    for item in book.get_items():
        if item.get_type() == 9:  # DOCUMENT
            soup = BeautifulSoup(item.get_content(), "html.parser")
            parts.append(soup.get_text(separator="\n", strip=True))
    return "\n".join(parts)

def extract_text_from_pages(file_path: str, max_pages: int, temp_dir: str) -> Tuple[str, Dict]:
    """
    Extract text from Apple Pages file (.pages).
    Pages files are zip archives containing QuickLook/Preview.pdf
    """
    import zipfile

    stats = {'pages_read': 0, 'ocr_used': False}

    try:
        with zipfile.ZipFile(file_path, 'r') as zip_ref:
            # Look for the preview PDF inside the .pages bundle
            preview_path = 'QuickLook/Preview.pdf'

            if preview_path not in zip_ref.namelist():
                # Some .pages files might have different structure
                # Try to find any PDF in the archive
                pdf_files = [name for name in zip_ref.namelist() if name.lower().endswith('.pdf')]
                if not pdf_files:
                    raise Exception("No PDF preview found in .pages file")
                preview_path = pdf_files[0]

            # Extract PDF to temp directory
            pdf_data = zip_ref.read(preview_path)
            temp_pdf = os.path.join(temp_dir, f"pages_preview_{uuid.uuid4()}.pdf")

            with open(temp_pdf, 'wb') as f:
                f.write(pdf_data)

            # Use existing PDF extraction
            text, pages_read = extract_text_from_pdf(temp_pdf, max_pages=max_pages)
            stats['pages_read'] = pages_read

            # Clean up temp PDF
            try:
                os.remove(temp_pdf)
            except Exception:
                pass

            return text, stats

    except zipfile.BadZipFile:
        raise Exception("Invalid .pages file (not a valid zip archive)")
    except Exception as e:
        raise Exception(f"Failed to extract from .pages file: {e}")

# --------------- Unified extractor ---------------

def extract_text_from_file(file_path: Path, max_pages: int, temp_dir: Optional[str],
                           max_rows_per_sheet: int, max_html_chars: int) -> Tuple[str, Dict]:
    stats = {'pages_read': 0, 'ocr_used': False}
    ext = file_path.suffix.lower()

    if ext == '.pdf':
        text, pages_read = extract_text_from_pdf(str(file_path), max_pages=max_pages)
        stats['pages_read'] = pages_read
        if not text.strip() and temp_dir is not None:
            ocr_path = _maybe_ocr_pdf(str(file_path), temp_dir)
            if ocr_path:
                stats['ocr_used'] = True
                text, pages_read = extract_text_from_pdf(ocr_path, max_pages=max_pages)
                stats['pages_read'] = pages_read
        return text, stats

    if ext == '.docx':
        return extract_text_from_docx(str(file_path)), stats

    if ext == '.doc':
        # Prefer conversion to DOCX for structured extraction
        if temp_dir:
            converted = _convert_with_soffice(str(file_path), "docx", temp_dir)
            if converted:
                return extract_text_from_docx(converted), stats
        # Fallback: antiword to plain text
        txt = _extract_with_antiword(str(file_path))
        if txt:
            return txt, stats
        raise Exception("Unable to process .doc (need LibreOffice or antiword)")

    if ext == '.pptx':
        return extract_text_from_pptx(str(file_path)), stats

    if ext == '.ppt':
        # Prefer conversion to PPTX for structured extraction
        if temp_dir:
            converted = _convert_with_soffice(str(file_path), "pptx", temp_dir)
            if converted:
                return extract_text_from_pptx(converted), stats
        # Fallback: catppt to plain text
        txt = _extract_with_catppt(str(file_path))
        if txt:
            return txt, stats
        raise Exception("Unable to process .ppt (need LibreOffice or catppt)")

    if ext == '.xlsx':
        return extract_text_from_xlsx(str(file_path), max_rows=max_rows_per_sheet), stats

    if ext == '.xls':
        return extract_text_from_xls(str(file_path), max_rows=max_rows_per_sheet), stats

    if ext == '.csv':
        return extract_text_from_csv(str(file_path), max_rows=max_rows_per_sheet), stats

    if ext in ('.html', '.htm'):
        return extract_text_from_html(str(file_path), max_chars=max_html_chars), stats

    if ext == '.rtf':
        return extract_text_from_rtf(str(file_path)), stats

    if ext == '.odt':
        return extract_text_from_odt(str(file_path)), stats

    if ext == '.epub':
        return extract_text_from_epub(str(file_path)), stats

    if ext == '.pages':
        if not temp_dir:
            raise Exception("Cannot process .pages file without temp directory")
        return extract_text_from_pages(str(file_path), max_pages=max_pages, temp_dir=temp_dir)

    if ext in ['.txt', '.md']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read(), stats

    raise Exception(f"Unsupported file type: {ext}")

# --------------- Chunk filtering & embeddings ---------------

def _validate_and_normalize_chunks(chunks: List[str], min_chunk_chars: int) -> List[str]:
    cleaned = []
    for c in chunks:
        s = re.sub(r'\s+', ' ', c).strip()
        if len(s) >= min_chunk_chars:
            cleaned.append(s)
    return cleaned

def _embed_single_batch(batch_texts: List[str], batch_num: int, openai_service) -> Tuple[int, List, int]:
    """
    Embed a single batch and return (batch_num, embeddings, tokens_used)
    Used for parallel embedding with ThreadPoolExecutor
    """
    try:
        resp = openai_service.create_embeddings_batch(batch_texts)
        if not resp or not resp.get('success'):
            error_msg = resp.get('error') if resp else 'unknown error'
            print(f"  Warning: embedding batch {batch_num} failed: {error_msg}")
            return (batch_num, [None] * len(batch_texts), 0)
        return (batch_num, resp['embeddings'], int(resp.get('tokens_used', 0)))
    except Exception as e:
        print(f"  Error in batch {batch_num}: {e}")
        return (batch_num, [None] * len(batch_texts), 0)


def _embed_in_batches(texts: List[str], embed_batch_size: int, max_workers: int = 8) -> Tuple[List[List[float]], int]:
    """
    Embed texts in parallel batches using ThreadPoolExecutor

    Args:
        texts: List of text chunks to embed
        embed_batch_size: Number of texts per API call
        max_workers: Number of parallel threads (default: 8)

    Returns:
        (embeddings, total_tokens)
    """
    openai_service = get_openai_service()

    # Split texts into batches
    batches = []
    for i in range(0, len(texts), embed_batch_size):
        batch = texts[i:i + embed_batch_size]
        batch_num = i // embed_batch_size + 1
        batches.append((batch_num, batch))

    total_batches = len(batches)
    print(f"  Embedding {len(texts)} chunks in {total_batches} batches using {max_workers} parallel workers...")

    # Process batches in parallel
    results = {}
    total_tokens = 0
    completed = 0

    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        # Submit all batches
        futures = {
            executor.submit(_embed_single_batch, batch_texts, batch_num, openai_service): batch_num
            for batch_num, batch_texts in batches
        }

        # Collect results as they complete
        for future in as_completed(futures):
            batch_num, embeddings, tokens_used = future.result()
            results[batch_num] = embeddings
            total_tokens += tokens_used
            completed += 1

            # Progress indicator
            if completed % max(1, total_batches // 10) == 0 or completed == total_batches:
                print(f"    Progress: {completed}/{total_batches} batches ({100*completed//total_batches}%)")

    # Reconstruct embeddings in original order
    all_embeddings = []
    for batch_num in range(1, total_batches + 1):
        all_embeddings.extend(results[batch_num])

    return all_embeddings, total_tokens

# --------------- File processor ---------------
SUPPORTED_EXTS = {'.csv', '.doc', '.docx', '.epub', '.htm', '.html', '.md',
                  '.odt', '.pages', '.pdf', '.ppt', '.pptx', '.rtf', '.txt', '.xls', '.xlsx'}

def process_file(
    file_path: Path,
    category: Optional[str],
    chunk_size: int,
    chunk_overlap: int,
    min_chunk_chars: int,
    text_snippet_len: int,
    embed_batch_size: int,
    max_chunks_per_file: int,
    max_pages: int,
    work_tmpdir: str,
    max_rows_per_sheet: int,
    max_html_chars: int,
    embedding_workers: int = 8,
) -> Optional[Dict]:
    print(f"Processing: {file_path.name}")

    ext = file_path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        print(f"  Skipping unsupported file type: {ext} (path: {file_path})")
        return None
    print(f"Processing: {file_path.name}")
    try:
        text, stats = extract_text_from_file(
            file_path=file_path,
            max_pages=max_pages,
            temp_dir=work_tmpdir,
            max_rows_per_sheet=max_rows_per_sheet,
            max_html_chars=max_html_chars,
        )
    except Exception as e:
        print(f"  Error extracting text: {e}")
        return None

    if not text or not text.strip():
        print("  No extractable text found. For scanned PDFs consider OCR:\n"
              "    sudo apt-get update && sudo apt-get install -y ocrmypdf tesseract-ocr")
        return None

    chunks = chunk_text(text, chunk_size=chunk_size, overlap=chunk_overlap)
    chunks = _validate_and_normalize_chunks(chunks, min_chunk_chars=min_chunk_chars)

    if not chunks:
        print("  No valid chunks after filtering. Skipping.")
        return None

    # Removed truncation - process all chunks regardless of count
    if max_chunks_per_file > 0 and len(chunks) > max_chunks_per_file:
        print(f"  ⚠️  Warning: {len(chunks)} chunks (exceeds --max-chunks-per-file={max_chunks_per_file})")
        print(f"     Processing ALL chunks with parallel embedding (no truncation)")
        # Note: Not truncating anymore - we'll parallelize embedding instead

    info = []
    if stats.get('pages_read'):
        info.append(f"pages read: {stats['pages_read']}")
    if stats.get('ocr_used'):
        info.append("OCR used")
    print(f"  Created {len(chunks)} valid chunks" + (f" | {', '.join(info)}" if info else ""))

    embeddings, tokens_used = _embed_in_batches(chunks, embed_batch_size=embed_batch_size, max_workers=embedding_workers)

    vectors, kept = [], 0
    for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
        if emb is None:
            continue
        metadata = {
            'text_snippet': chunk[:text_snippet_len],
            'full_text': chunk,  # Store full text for BM25 indexing
            'source': file_path.name,
            'file_path': str(file_path),
            'file_type': file_path.suffix[1:],
            'chunk_index': i,
            'total_chunks': len(chunks),
            'ingested_at': time.strftime('%Y-%m-%d %H:%M:%S')
        }
        if category:
            metadata['category'] = category
        vectors.append({'id': str(uuid.uuid4()), 'values': emb, 'metadata': metadata})
        kept += 1

    print(f"  Prepared {kept} vectors")
    if kept == 0:
        print("  All embedding batches failed; skipping file.")
        return None

    return {'vectors': vectors, 'tokens_used': tokens_used}

# --------------- S3 helpers ---------------

def list_s3_objects(s3_client, bucket: str, prefix: str = "") -> List[Dict]:
    """
    List all objects under a prefix. Do NOT filter by extension here.
    We'll filter at processing time. This avoids 'No supported files found'
    when filenames are odd or slightly malformed.
    """
    objects = []
    paginator = s3_client.get_paginator('list_objects_v2')

    # Debugging counters
    total_items = 0
    skipped_dirs = []
    skipped_empty = []

    try:
        for page in paginator.paginate(Bucket=bucket, Prefix=prefix):
            for obj in page.get('Contents', []):
                key = obj['Key']
                size = obj.get('Size', 0)
                total_items += 1

                # Skip "folder markers" (keys that end with '/')
                if key.endswith('/'):
                    skipped_dirs.append(key)
                    continue

                if size == 0:
                    skipped_empty.append(key)
                    continue

                objects.append({
                    'key': key,
                    'size': size,
                    'last_modified': obj.get('LastModified')
                })
    except ClientError as e:
        print(f"Error listing S3 objects: {e}")
        return []

    # Enhanced debug output
    print(f"\n=== S3 Listing Debug ===")
    print(f"Total items found: {total_items}")
    print(f"Skipped (directories ending with '/'): {len(skipped_dirs)}")
    if skipped_dirs[:5]:
        for d in skipped_dirs[:5]:
            print(f"  - {d}")
        if len(skipped_dirs) > 5:
            print(f"  ... and {len(skipped_dirs) - 5} more")

    print(f"Skipped (zero-byte files): {len(skipped_empty)}")
    if skipped_empty[:5]:
        for e in skipped_empty[:5]:
            print(f"  - {e}")
        if len(skipped_empty) > 5:
            print(f"  ... and {len(skipped_empty) - 5} more")

    print(f"Files to process: {len(objects)}")
    if objects:
        print(f"Sample files (first 10):")
        for sample in objects[:10]:
            ext = Path(sample['key']).suffix.lower()
            supported = "✓" if ext in SUPPORTED_EXTS else "✗"
            print(f"  {supported} {sample['key']} ({sample['size']:,} bytes, ext={ext})")
        if len(objects) > 10:
            print(f"  ... and {len(objects) - 10} more files")
    print(f"======================\n")

    return objects

# --------------- Drive helpers (updated exports) ---------------

def list_gdrive_files(service, folder_id: str = None) -> List[Dict]:
    files, page_token = [], None
    try:
        q = ["trashed = false"]
        if folder_id:
            q.append(f"'{folder_id}' in parents")
        query = " and ".join(q)
        while True:
            results = service.files().list(
                q=query,
                spaces='drive',
                fields='nextPageToken, files(id, name, mimeType, modifiedTime, size)',
                pageToken=page_token,
                pageSize=100
            ).execute()
            for item in results.get('files', []):
                files.append(item)
            page_token = results.get('nextPageToken')
            if not page_token:
                break
    except Exception as e:
        print(f"Error listing Google Drive files: {e}")
        return []
    return files

def download_gdrive_file(service, file_id: str, file_name: str, mime_type: str, local_path: str) -> bool:
    try:
        # Google Docs → DOCX
        if mime_type == 'application/vnd.google-apps.document':
            request = service.files().export_media(
                fileId=file_id,
                mimeType='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            )
            local_path = str(Path(local_path).with_suffix('.docx'))

        # Sheets → CSV
        elif mime_type == 'application/vnd.google-apps.spreadsheet':
            request = service.files().export_media(
                fileId=file_id,
                mimeType='text/csv'
            )
            local_path = str(Path(local_path).with_suffix('.csv'))

        # Slides → PPTX
        elif mime_type == 'application/vnd.google-apps.presentation':
            request = service.files().export_media(
                fileId=file_id,
                mimeType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            local_path = str(Path(local_path).with_suffix('.pptx'))

        # Other Google apps formats: skip
        elif mime_type.startswith('application/vnd.google-apps'):
            print(f"  Skipping unsupported Google Apps file: {file_name} ({mime_type})")
            return False

        else:
            # Regular download
            request = service.files().get_media(fileId=file_id)

        with io.FileIO(local_path, 'wb') as fh:
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
        return True
    except Exception as e:
        print(f"Error downloading {file_name}: {e}")
        return False

# --------------- Streaming uploader ---------------

def _flush_buffer(pinecone_service, namespace: str, buffer: List[Dict], uploaded_count: int, bm25_service=None) -> int:
    """
    Flush vector buffer to Pinecone and optionally update BM25 index

    Args:
        pinecone_service: Pinecone service instance
        namespace: Namespace to upload to
        buffer: List of vector dictionaries
        uploaded_count: Running count of uploaded vectors
        bm25_service: Optional BM25 service instance for keyword indexing

    Returns:
        Updated uploaded_count
    """
    if not buffer:
        return uploaded_count

    # Upload to Pinecone
    resp = pinecone_service.upsert_vectors(namespace, buffer)
    if resp.get('success'):
        count = resp.get('upserted_count', 0)
        uploaded_count += count
        print(f"  Uploaded batch: {count} vectors (total uploaded: {uploaded_count})")

        # Update BM25 index if service provided
        if bm25_service:
            try:
                # Extract documents and IDs for BM25
                doc_texts = []
                doc_ids = []
                for vec in buffer:
                    # Get full_text from metadata
                    full_text = vec.get('metadata', {}).get('full_text', '')
                    if full_text:
                        doc_texts.append(full_text)
                        doc_ids.append(vec['id'])

                if doc_texts:
                    # Add to BM25 index (this will rebuild the index for the namespace)
                    bm25_result = bm25_service.add_documents(namespace, doc_texts, doc_ids, append=True)
                    if bm25_result['success']:
                        print(f"  BM25 index updated: {bm25_result['document_count']} documents")
                    else:
                        print(f"  Warning: BM25 indexing failed: {bm25_result.get('error')}")
            except Exception as e:
                print(f"  Warning: BM25 indexing error: {e}")
    else:
        print(f"  Error uploading batch: {resp.get('error')}")

    buffer.clear()
    return uploaded_count

# --------------- Ingestors ---------------
def download_s3_file(s3_client, bucket: str, key: str, local_path: str) -> bool:
    try:
        s3_client.download_file(bucket, key, local_path)
        return True
    except ClientError as e:
        print(f"Error downloading {key}: {e}")
        return False

def ingest_directory(directory: str, namespace: str, category: Optional[str], batch_size: int,
                     text_snippet_len: int, chunk_size: int, chunk_overlap: int,
                     min_chunk_chars: int, embed_batch_size: int, max_chunks_per_file: int, max_pages: int,
                     max_rows_per_sheet: int, max_html_chars: int, embedding_workers: int = 8):
    directory_path = Path(directory)
    if not directory_path.exists():
        print(f"Error: Directory not found: {directory}")
        return

    files: List[Path] = []
    for ext in ALLOWED_EXTENSIONS:
        files.extend(directory_path.rglob(f'*{ext}'))

    if not files:
        print(f"No supported files found in {directory}")
        print(f"Supported extensions: {', '.join(sorted(ALLOWED_EXTENSIONS))}")
        return

    print(f"\nFound {len(files)} files to process")
    print(f"Target namespace: {namespace}")
    print("=" * 60)

    pinecone_service = get_pinecone_service()
    bm25_service = get_bm25_service()
    buffer: List[Dict] = []
    uploaded_count = 0
    total_tokens = 0
    successful_files = 0
    failed_files = 0

    with tempfile.TemporaryDirectory(prefix='ingest_local_') as work_tmpdir:
        for file_path in files:
            result = process_file(
                file_path=file_path,
                category=category,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                min_chunk_chars=min_chunk_chars,
                text_snippet_len=text_snippet_len,
                embed_batch_size=embed_batch_size,
                max_chunks_per_file=max_chunks_per_file,
                max_pages=max_pages,
                work_tmpdir=work_tmpdir,
                max_rows_per_sheet=max_rows_per_sheet,
                max_html_chars=max_html_chars,
                embedding_workers=embedding_workers,
            )
            if result:
                buffer.extend(result['vectors'])
                total_tokens += result['tokens_used']
                successful_files += 1
                while len(buffer) >= batch_size:
                    uploaded_count = _flush_buffer(pinecone_service, namespace, buffer[:batch_size], uploaded_count, bm25_service)
                    del buffer[:batch_size]
            else:
                failed_files += 1

    uploaded_count = _flush_buffer(pinecone_service, namespace, buffer, uploaded_count, bm25_service)

    # Save BM25 index to storage (S3 or local)
    if bm25_service and namespace in bm25_service.get_all_namespaces():
        print("\nSaving BM25 index to storage...")
        save_result = bm25_service.save_index(namespace)
        if save_result['success']:
            storage = save_result.get('storage', 'unknown')
            location = save_result.get('s3_key') or save_result.get('filepath', 'unknown')
            print(f"  ✓ BM25 index saved to {storage}: {location}")
        else:
            print(f"  ✗ Failed to save BM25 index: {save_result.get('error')}")

    print("\n" + "=" * 60)
    print("Processing complete:")
    print(f"  Successful files: {successful_files}")
    print(f"  Failed files:     {failed_files}")
    print(f"  Total tokens:     {total_tokens:,}")
    print(f"  Estimated cost:   ${(total_tokens * 0.0001 / 1000):.4f}")
    print(f"  Total uploaded:   {uploaded_count} vectors")

def ingest_from_s3(bucket: str, prefix: str, namespace: str, category: Optional[str], batch_size: int,
                   text_snippet_len: int, chunk_size: int, chunk_overlap: int,
                   min_chunk_chars: int, embed_batch_size: int, max_chunks_per_file: int, max_pages: int,
                   region: Optional[str], max_rows_per_sheet: int, max_html_chars: int, embedding_workers: int = 8):
    if not S3_SUPPORT:
        print("Error: boto3 not installed. Run: pip install boto3")
        return
    try:
        s3_client = boto3.client('s3', region_name=region) if region else boto3.client('s3')
    except NoCredentialsError:
        print("Error: AWS credentials not found. Configure via env vars, ~/.aws/credentials, or an IAM role.")
        return

    print(f"\nListing files in s3://{bucket}/{prefix}")
    objects = list_s3_objects(s3_client, bucket, prefix)
    if not objects:
        print(f"No supported files found in s3://{bucket}/{prefix}")
        print(f"Supported extensions: {', '.join(sorted(ALLOWED_EXTENSIONS))}")
        return

    print(f"\nFound {len(objects)} files to process")
    print(f"Target namespace: {namespace}")
    print("=" * 60)

    pinecone_service = get_pinecone_service()
    bm25_service = get_bm25_service()
    buffer: List[Dict] = []
    uploaded_count = 0
    total_tokens = 0
    successful_files = 0
    failed_files = 0

    temp_dir = tempfile.mkdtemp(prefix='s3_ingest_')
    try:
        for obj in objects:
            s3_key = obj['key']
            filename = Path(s3_key).name
            local_path = os.path.join(temp_dir, filename)

            print(f"\nDownloading: {s3_key}")
            if not download_s3_file(s3_client, bucket, s3_key, local_path):
                failed_files += 1
                continue

            result = process_file(
                file_path=Path(local_path),
                category=category,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                min_chunk_chars=min_chunk_chars,
                text_snippet_len=text_snippet_len,
                embed_batch_size=embed_batch_size,
                max_chunks_per_file=max_chunks_per_file,
                max_pages=max_pages,
                work_tmpdir=temp_dir,
                max_rows_per_sheet=max_rows_per_sheet,
                max_html_chars=max_html_chars,
                embedding_workers=embedding_workers,
            )

            try:
                os.remove(local_path)
            except Exception:
                pass

            if result:
                buffer.extend(result['vectors'])
                total_tokens += result['tokens_used']
                successful_files += 1
                while len(buffer) >= batch_size:
                    uploaded_count = _flush_buffer(pinecone_service, namespace, buffer[:batch_size], uploaded_count, bm25_service)
                    del buffer[:batch_size]
            else:
                failed_files += 1
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    uploaded_count = _flush_buffer(pinecone_service, namespace, buffer, uploaded_count, bm25_service)

    # Save BM25 index to storage (S3 or local)
    if bm25_service and namespace in bm25_service.get_all_namespaces():
        print("\nSaving BM25 index to storage...")
        save_result = bm25_service.save_index(namespace)
        if save_result['success']:
            storage = save_result.get('storage', 'unknown')
            location = save_result.get('s3_key') or save_result.get('filepath', 'unknown')
            print(f"  ✓ BM25 index saved to {storage}: {location}")
        else:
            print(f"  ✗ Failed to save BM25 index: {save_result.get('error')}")

    print("\n" + "=" * 60)
    print("Processing complete:")
    print(f"  Successful files: {successful_files}")
    print(f"  Failed files:     {failed_files}")
    print(f"  Total tokens:     {total_tokens:,}")
    print(f"  Estimated cost:   ${(total_tokens * 0.0001 / 1000):.4f}")
    print(f"  Total uploaded:   {uploaded_count} vectors")

def ingest_from_gdrive(folder_id: Optional[str], credentials_path: str, namespace: str, category: Optional[str], batch_size: int,
                       text_snippet_len: int, chunk_size: int, chunk_overlap: int,
                       min_chunk_chars: int, embed_batch_size: int, max_chunks_per_file: int, max_pages: int,
                       max_rows_per_sheet: int, max_html_chars: int, embedding_workers: int = 8):
    if not GDRIVE_SUPPORT:
        print("Error: Google Drive libraries not installed.")
        print("Run: pip install google-auth google-auth-oauthlib google-auth-httplib2 google-api-python-client")
        return
    try:
        SCOPES = ['https://www.googleapis.com/auth/drive.readonly']
        credentials = service_account.Credentials.from_service_account_file(credentials_path, scopes=SCOPES)
        service = build('drive', 'v3', credentials=credentials)
    except Exception as e:
        print(f"Error initializing Google Drive API: {e}")
        print("Ensure your service account credentials JSON is valid.")
        return

    print(f"\nListing files in Google Drive folder: {folder_id or 'root'}")
    files = list_gdrive_files(service, folder_id)
    if not files:
        print("No supported files found in Google Drive folder")
        print(f"Supported extensions: {', '.join(sorted(ALLOWED_EXTENSIONS))}")
        return

    print(f"\nFound {len(files)} files to process")
    print(f"Target namespace: {namespace}")
    print("=" * 60)

    pinecone_service = get_pinecone_service()
    bm25_service = get_bm25_service()
    buffer: List[Dict] = []
    uploaded_count = 0
    total_tokens = 0
    successful_files = 0
    failed_files = 0

    temp_dir = tempfile.mkdtemp(prefix='gdrive_ingest_')
    try:
        for info in files:
            file_id = info['id']
            file_name = info['name']
            mime_type = info.get('mimeType', 'application/octet-stream')  # Fixed: API returns 'mimeType', not 'mime_type'
            local_path = os.path.join(temp_dir, file_name)

            print(f"\nDownloading: {file_name}")
            if not download_gdrive_file(service, file_id, file_name, mime_type, local_path):
                failed_files += 1
                continue

            result = process_file(
                file_path=Path(local_path),
                category=category,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                min_chunk_chars=min_chunk_chars,
                text_snippet_len=text_snippet_len,
                embed_batch_size=embed_batch_size,
                max_chunks_per_file=max_chunks_per_file,
                max_pages=max_pages,
                work_tmpdir=temp_dir,
                max_rows_per_sheet=max_rows_per_sheet,
                max_html_chars=max_html_chars,
                embedding_workers=embedding_workers,
            )

            try:
                if os.path.exists(local_path):
                    os.remove(local_path)
            except Exception:
                pass

            if result:
                buffer.extend(result['vectors'])
                total_tokens += result['tokens_used']
                successful_files += 1
                while len(buffer) >= batch_size:
                    uploaded_count = _flush_buffer(pinecone_service, namespace, buffer[:batch_size], uploaded_count, bm25_service)
                    del buffer[:batch_size]
            else:
                failed_files += 1
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

    uploaded_count = _flush_buffer(pinecone_service, namespace, buffer, uploaded_count, bm25_service)

    # Save BM25 index to storage (S3 or local)
    if bm25_service and namespace in bm25_service.get_all_namespaces():
        print("\nSaving BM25 index to storage...")
        save_result = bm25_service.save_index(namespace)
        if save_result['success']:
            storage = save_result.get('storage', 'unknown')
            location = save_result.get('s3_key') or save_result.get('filepath', 'unknown')
            print(f"  ✓ BM25 index saved to {storage}: {location}")
        else:
            print(f"  ✗ Failed to save BM25 index: {save_result.get('error')}")

    print("\n" + "=" * 60)
    print("Processing complete:")
    print(f"  Successful files: {successful_files}")
    print(f"  Failed files:     {failed_files}")
    print(f"  Total tokens:     {total_tokens:,}")
    print(f"  Estimated cost:   ${(total_tokens * 0.0001 / 1000):.4f}")
    print(f"  Total uploaded:   {uploaded_count} vectors")

# --------------- CLI ---------------

def main():
    parser = argparse.ArgumentParser(
        description='Bulk ingest documents from local directory, AWS S3, or Google Drive into Pinecone (streaming uploads, safe batching, multi-format incl. .doc/.ppt)',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  LOCAL DIRECTORY:
    python scripts/bulk_ingest.py /path/to/docs --namespace shared --batch-size 25

  AWS S3:
    python scripts/bulk_ingest.py --s3-bucket my-bucket --s3-prefix "documents/policies/" --namespace shared --batch-size 25

  GOOGLE DRIVE:
    python scripts/bulk_ingest.py --gdrive-folder 1abc123xyz --gdrive-credentials creds.json --namespace shared --batch-size 25

Notes:
  • Use smaller --batch-size for Pinecone upserts if memory is constrained (e.g., 25).
  • --embed-batch-size controls OpenAI embedding micro-batch size (default 64).
  • --max-pages limits PDF pages read (helps with massive books).
  • --max-chunks-per-file caps per-file chunk count.
  • Spreadsheet row cap (--max-rows-per-sheet) and HTML size cap (--max-html-chars).
  • Legacy Office support: .doc/.ppt via LibreOffice (preferred) or antiword/catppt fallbacks.
        """
    )

    parser.add_argument('path', nargs='?', help='Local directory containing documents (if not using S3 or Google Drive)')
    parser.add_argument('--namespace', required=True, help='Pinecone namespace to ingest into')
    parser.add_argument('--category', default=None, help='Optional category tag for all documents')

    # Streaming upsert buffer size for Pinecone
    parser.add_argument('--batch-size', type=int, default=100, help='Vectors per Pinecone upsert batch (default: 100)')

    # Embedding + chunk config
    parser.add_argument('--chunk-size', type=int, default=DEFAULT_CHUNK_SIZE, help='Chunk size in characters (default: 1000)')
    parser.add_argument('--chunk-overlap', type=int, default=DEFAULT_CHUNK_OVERLAP, help='Chunk overlap in characters (default: 200)')
    parser.add_argument('--min-chunk-chars', type=int, default=25, help='Minimum chars for a chunk to be embedded (default: 25)')
    parser.add_argument('--text-snippet-len', type=int, default=500, help='Chars of chunk text to store in metadata (default: 500)')
    parser.add_argument('--embed-batch-size', type=int, default=64, help='Texts per embeddings API call (default: 64)')
    parser.add_argument('--embedding-workers', type=int, default=8, help='Number of parallel threads for embedding (default: 8, max recommended: 16)')
    parser.add_argument('--max-chunks-per-file', type=int, default=0, help='DEPRECATED: Truncation removed. All chunks are now processed with parallel embedding.')
    parser.add_argument('--max-pages', type=int, default=0, help='Only read first N pages of PDFs (0 = all pages)')

    # New caps for other formats
    parser.add_argument('--max-rows-per-sheet', type=int, default=5000, help='Max rows per sheet for CSV/XLS/XLSX (default: 5000)')
    parser.add_argument('--max-html-chars', type=int, default=2_000_000, help='Max characters to keep from HTML (default: 2,000,000)')

    # S3 options
    parser.add_argument('--s3-bucket', help='AWS S3 bucket name')
    parser.add_argument('--s3-prefix', default='', help='S3 prefix/path (e.g., "documents/policies/")')
    parser.add_argument('--s3-region', help='AWS region (default: from AWS config)')

    # Google Drive options
    parser.add_argument('--gdrive-folder', help='Google Drive folder ID (omit for root folder)')
    parser.add_argument('--gdrive-credentials', help='Path to Google service account credentials JSON file')

    args = parser.parse_args()

    print("Bulk Document Ingestion (streaming + batching + multi-format incl. .doc/.ppt)")
    print("=" * 60)

    if args.s3_bucket:
        ingest_from_s3(
            bucket=args.s3_bucket,
            prefix=args.s3_prefix,
            namespace=args.namespace,
            category=args.category,
            batch_size=args.batch_size,
            text_snippet_len=args.text_snippet_len,
            chunk_size=args.chunk_size,
            chunk_overlap=args.chunk_overlap,
            min_chunk_chars=args.min_chunk_chars,
            embed_batch_size=args.embed_batch_size,
            max_chunks_per_file=args.max_chunks_per_file,
            max_pages=args.max_pages,
            region=args.s3_region,
            max_rows_per_sheet=args.max_rows_per_sheet,
            max_html_chars=args.max_html_chars,
            embedding_workers=args.embedding_workers,
        )
    elif args.gdrive_credentials:
        ingest_from_gdrive(
            folder_id=args.gdrive_folder,
            credentials_path=args.gdrive_credentials,
            namespace=args.namespace,
            category=args.category,
            batch_size=args.batch_size,
            text_snippet_len=args.text_snippet_len,
            chunk_size=args.chunk_size,
            chunk_overlap=args.chunk_overlap,
            min_chunk_chars=args.min_chunk_chars,
            embed_batch_size=args.embed_batch_size,
            max_chunks_per_file=args.max_chunks_per_file,
            max_pages=args.max_pages,
            max_rows_per_sheet=args.max_rows_per_sheet,
            max_html_chars=args.max_html_chars,
            embedding_workers=args.embedding_workers,
        )
    elif args.path:
        ingest_directory(
            directory=args.path,
            namespace=args.namespace,
            category=args.category,
            batch_size=args.batch_size,
            text_snippet_len=args.text_snippet_len,
            chunk_size=args.chunk_size,
            chunk_overlap=args.chunk_overlap,
            min_chunk_chars=args.min_chunk_chars,
            embed_batch_size=args.embed_batch_size,
            max_chunks_per_file=args.max_chunks_per_file,
            max_pages=args.max_pages,
            max_rows_per_sheet=args.max_rows_per_sheet,
            max_html_chars=args.max_html_chars,
            embedding_workers=args.embedding_workers,
        )
    else:
        print("Error: You must specify either:")
        print("  - A local directory path")
        print("  - --s3-bucket for S3 ingestion")
        print("  - --gdrive-credentials for Google Drive ingestion")
        print("\nUse --help for more information")
        sys.exit(1)

if __name__ == '__main__':
    main()
